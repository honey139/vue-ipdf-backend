import sys
import os
from pptx import Presentation
from pdf2image import convert_from_path
from pptx.util import Inches, Pt

def convert_pdf_to_pptx(pdf_path, pptx_path, resolution=150):
    try:
        print(f"Starting conversion: {pdf_path} to {pptx_path}")
        sys.stdout.flush()  # Ensure messages are sent to Node.js or other callers

        # Ensure the input file exists
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"Input file does not exist: {pdf_path}")

        # Convert all PDF pages to images
        images = convert_from_path(pdf_path, dpi=resolution)

        if not images:
            print(f"No pages found in the PDF: {pdf_path}")
            sys.stdout.flush()
            return False

        # Create a PowerPoint presentation object
        presentation = Presentation()

        # Get the first page dimensions for setting slide size
        first_image = images[0]
        pdf_width, pdf_height = first_image.size  # Width and height in pixels

        # Convert pixel dimensions to PowerPoint units (points)
        # 1 inch = 72 points, 1 inch = 96 DPI
        ppt_width = Pt(pdf_width * 72 / resolution)
        ppt_height = Pt(pdf_height * 72 / resolution)

        # Set the slide dimensions
        presentation.slide_width = ppt_width
        presentation.slide_height = ppt_height

        for i, image in enumerate(images):
            # Save the image temporarily
            temp_image_path = f"temp_image_{i}.png"
            image.save(temp_image_path, optimize=True, quality=85)

            # Add a blank slide to the presentation
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Layout 6 = blank slide
            
            # Insert the image into the slide
            slide.shapes.add_picture(temp_image_path, 0, 0, width=ppt_width, height=ppt_height)

            # Remove the temporary image file
            os.remove(temp_image_path)

        # Save the presentation to the specified path
        presentation.save(pptx_path)

        if os.path.exists(pptx_path):
            print(f"Conversion successful: {pptx_path}")
            sys.stdout.flush()
            return True
        else:
            print(f"Conversion failed: Output file not created at {pptx_path}")
            sys.stdout.flush()
            return False

    except FileNotFoundError as e:
        print(f"File not found: {e}")
        sys.stdout.flush()
        return False
    except ImportError as e:
        print(f"Import error: {e}")
        sys.stdout.flush()
        return False
    except Exception as e:
        print(f"Unexpected error: {e}")
        sys.stdout.flush()
        return False

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python convertPdfToPpt.py <pdf_path> <pptx_path>")
        sys.exit(1)

    pdf_path = sys.argv[1]
    pptx_path = sys.argv[2]

    success = convert_pdf_to_pptx(pdf_path, pptx_path, resolution=150)
    sys.exit(0 if success else 1)
